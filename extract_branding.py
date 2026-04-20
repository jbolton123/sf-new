from pptx import Presentation
from pptx.util import Pt
from pptx.enum.dml import MSO_THEME_COLOR
import os

def extract_branding(filepath, output_name):
    try:
        prs = Presentation(filepath)
        print(f"\n=== {output_name} ===")
        print(f"Dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
        print(f"Total slides: {len(prs.slides)}")

        colors_found = set()
        fonts_found = set()

        # Check first 5 slides for colors and fonts
        for i, slide in enumerate(prs.slides[:5]):
            for shape in slide.shapes:
                # Extract text formatting
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                # Font
                                if run.font.name:
                                    fonts_found.add(run.font.name)
                                # Color
                                if hasattr(run.font.color, 'rgb'):
                                    colors_found.add(str(run.font.color.rgb))
                            except:
                                pass

                # Extract shape fill colors
                if hasattr(shape, "fill"):
                    try:
                        if hasattr(shape.fill.fore_color, 'rgb'):
                            colors_found.add(str(shape.fill.fore_color.rgb))
                    except:
                        pass

        print(f"\nFonts used: {', '.join(sorted(fonts_found)) if fonts_found else 'None found'}")
        print(f"Colors found: {len(colors_found)}")
        for color in sorted(colors_found)[:10]:
            print(f"  - {color}")

        return True
    except Exception as e:
        print(f"Error: {e}")
        return False

files = [
    ("/Users/j.bolton/Downloads/FY26 Loyalty Management Value Book.pptx", "FY26 Value Book"),
    ("/Users/j.bolton/Downloads/Jack Presentation.pptx", "Jack Presentation"),
    ("/Users/j.bolton/Downloads/SpaceNK-NDULGE-Customer-Deck.pptx", "SpaceNK Customer Deck"),
]

for filepath, name in files:
    if os.path.exists(filepath):
        extract_branding(filepath, name)
