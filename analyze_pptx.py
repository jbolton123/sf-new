from pptx import Presentation
import sys

def analyze_presentation(filepath):
    prs = Presentation(filepath)

    print(f"\n=== Analyzing: {filepath.split('/')[-1]} ===\n")
    print(f"Slide dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
    print(f"Total slides: {len(prs.slides)}")

    # Analyze first 3 slides
    for i, slide in enumerate(prs.slides[:3]):
        print(f"\n--- Slide {i+1} ---")
        print(f"Layout: {slide.slide_layout.name}")

        # Check for images (potential logos)
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                print(f"  Image found: {shape.width.inches}\" x {shape.height.inches}\"")

            # Get text content
            if hasattr(shape, "text") and shape.text.strip():
                text_preview = shape.text[:100].replace('\n', ' ')
                print(f"  Text: {text_preview}")

            # Get fill colors
            if hasattr(shape, "fill"):
                try:
                    if shape.fill.type == 1:  # Solid fill
                        rgb = shape.fill.fore_color.rgb
                        print(f"  Fill color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")
                except:
                    pass

            # Get text colors
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            if hasattr(run.font, "color") and run.font.color.type == 1:
                                rgb = run.font.color.rgb
                                print(f"  Text color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")
                                if run.font.name:
                                    print(f"  Font: {run.font.name}, Size: {run.font.size.pt if run.font.size else 'default'}pt")
                                break
                        except:
                            pass

if __name__ == "__main__":
    files = [
        "/Users/j.bolton/Downloads/FY26 Loyalty Management Value Book.pptx",
        "/Users/j.bolton/Downloads/Jack Presentation.pptx",
        "/Users/j.bolton/Downloads/FY27 Next Gen Marketing Cloud 201 Product Features Deck.pptx"
    ]

    for f in files:
        try:
            analyze_presentation(f)
        except Exception as e:
            print(f"Error analyzing {f}: {e}")
