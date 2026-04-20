from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Q1 2026 Results"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Strong Growth Across All Metrics"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

# Slide 2: Key Highlights
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf2 = title2.text_frame
tf2.text = "Key Highlights"
tf2.paragraphs[0].font.size = Pt(40)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content2 = slide2.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
tf2c = content2.text_frame
highlights = [
    "Revenue increased 42% YoY to $8.7M",
    "Customer acquisition up 67% - added 234 new customers",
    "Product launch exceeded targets by 28%",
    "Team expansion: 12 new hires across engineering and sales",
    "Net Promoter Score improved from 68 to 79"
]
for i, highlight in enumerate(highlights):
    if i == 0:
        tf2c.text = "• " + highlight
        p = tf2c.paragraphs[0]
    else:
        p = tf2c.add_paragraph()
        p.text = "• " + highlight
    p.font.size = Pt(24)
    p.space_before = Pt(15)
    p.level = 0

# Slide 3: Revenue Breakdown
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf3 = title3.text_frame
tf3.text = "Revenue Breakdown"
tf3.paragraphs[0].font.size = Pt(40)
tf3.paragraphs[0].font.bold = True
tf3.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content3 = slide3.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
tf3c = content3.text_frame
revenue_data = [
    "Total Revenue: $8.7M (+42% YoY)",
    "",
    "By Product Line:",
    "  → Enterprise: $5.2M (60%)",
    "  → SMB: $2.4M (28%)",
    "  → Self-Service: $1.1M (12%)",
    "",
    "By Region:",
    "  → North America: $6.1M (70%)",
    "  → EMEA: $1.7M (20%)",
    "  → APAC: $0.9M (10%)"
]
for i, line in enumerate(revenue_data):
    if i == 0:
        tf3c.text = line
        p = tf3c.paragraphs[0]
    else:
        p = tf3c.add_paragraph()
        p.text = line
    p.font.size = Pt(20)
    p.space_before = Pt(8)

# Slide 4: Customer Growth
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf4 = title4.text_frame
tf4.text = "Customer Growth & Retention"
tf4.paragraphs[0].font.size = Pt(40)
tf4.paragraphs[0].font.bold = True
tf4.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content4 = slide4.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
tf4c = content4.text_frame
customer_data = [
    "New Customers: 234 (+67% vs Q1 2025)",
    "Total Active Customers: 1,456",
    "Churn Rate: 2.3% (industry avg: 5.1%)",
    "Customer Lifetime Value: $47,300",
    "Average Deal Size: $37,200 (+23% QoQ)",
    "",
    "Top 3 Customer Wins:",
    "  • Global Financial Services - $420K ARR",
    "  • Healthcare Platform - $280K ARR",
    "  • E-commerce Leader - $195K ARR"
]
for i, line in enumerate(customer_data):
    if i == 0:
        tf4c.text = line
        p = tf4c.paragraphs[0]
    else:
        p = tf4c.add_paragraph()
        p.text = line
    p.font.size = Pt(20)
    p.space_before = Pt(10)

# Slide 5: Q2 Outlook
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf5 = title5.text_frame
tf5.text = "Q2 2026 Outlook"
tf5.paragraphs[0].font.size = Pt(40)
tf5.paragraphs[0].font.bold = True
tf5.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content5 = slide5.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
tf5c = content5.text_frame
outlook = [
    "Revenue Target: $10.2M (+17% QoQ)",
    "New Customer Target: 180-200",
    "",
    "Key Initiatives:",
    "  • Launch AI-powered analytics module (May)",
    "  • Expand into 3 new verticals",
    "  • Open APAC office in Singapore",
    "  • Release mobile app (beta in June)",
    "",
    "Strategic Focus:",
    "  → Accelerate enterprise pipeline",
    "  → Enhance product-market fit in healthcare",
    "  → Scale customer success operations"
]
for i, line in enumerate(outlook):
    if i == 0:
        tf5c.text = line
        p = tf5c.paragraphs[0]
    else:
        p = tf5c.add_paragraph()
        p.text = line
    p.font.size = Pt(20)
    p.space_before = Pt(10)

# Save presentation
prs.save('Q1_2026_Results.pptx')
print("Presentation created successfully: Q1_2026_Results.pptx")
