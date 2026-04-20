from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Salesforce 2025 Corporate Brand Colors (from theme1.xml)
SF_DARK_BLUE = RGBColor(3, 45, 96)      # #032D60 - Primary dark
SF_LIGHT_BLUE = RGBColor(13, 157, 218)  # #0D9DDA - Accent/highlight
SF_MED_BLUE = RGBColor(11, 92, 171)     # #0B5CAB - Accent 1
SF_TEAL = RGBColor(11, 130, 124)        # #0B827C - Accent 2
SF_PINK = RGBColor(182, 5, 84)          # #B60554 - Accent 3
SF_ORANGE = RGBColor(221, 122, 1)       # #DD7A01 - Accent 4
SF_LIGHT = RGBColor(144, 208, 254)      # #90D0FE - Accent 5
SF_PURPLE = RGBColor(72, 26, 84)        # #481A54 - Accent 6
SF_GRAY = RGBColor(68, 68, 68)          # #444444
WHITE = RGBColor(255, 255, 255)

# Create presentation with proper 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

def add_footer(slide, page_num=None):
    """Add Salesforce footer bar to slide"""
    footer_height = Inches(0.35)
    footer = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0),
        prs.slide_height - footer_height,
        prs.slide_width,
        footer_height
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = SF_DARK_BLUE
    footer.line.fill.background()

    if page_num:
        page_box = slide.shapes.add_textbox(
            prs.slide_width - Inches(0.8),
            prs.slide_height - footer_height + Inches(0.08),
            Inches(0.5),
            Inches(0.2)
        )
        tf = page_box.text_frame
        tf.text = str(page_num)
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.RIGHT

def add_header_bar(slide):
    """Add decorative header element"""
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.05)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = SF_LIGHT_BLUE
    header.line.fill.background()

# ============= SLIDE 1: TITLE SLIDE =============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Background gradient effect (simplified as solid with decorative elements)
bg = slide1.shapes.add_shape(
    1,  # Rectangle
    Inches(0), Inches(0),
    prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = SF_DARK_BLUE
bg.line.fill.background()

# Accent bar
accent = slide1.shapes.add_shape(
    1,
    Inches(0), Inches(2.2),
    Inches(0.15), Inches(1.2)
)
accent.fill.solid()
accent.fill.fore_color.rgb = SF_LIGHT_BLUE
accent.line.fill.background()

# Main title
title = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
tf = title.text_frame
tf.text = "Q1 2026"
p = tf.paragraphs[0]
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Arial"

# Subtitle with Salesforce brand styling
subtitle = slide1.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(0.8))
tf_sub = subtitle.text_frame
tf_sub.text = "RESULTS"
p_sub = tf_sub.paragraphs[0]
p_sub.font.size = Pt(48)
p_sub.font.bold = True
p_sub.font.color.rgb = SF_LIGHT_BLUE

# Tagline
tag = slide1.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.4))
tf_tag = tag.text_frame
tf_tag.text = "Strong Growth Across All Metrics"
p_tag = tf_tag.paragraphs[0]
p_tag.font.size = Pt(20)
p_tag.font.color.rgb = SF_LIGHT
p_tag.font.name = "Arial"

# Date
date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.3))
tf_date = date_box.text_frame
tf_date.text = "April 2026"
p_date = tf_date.paragraphs[0]
p_date.font.size = Pt(14)
p_date.font.color.rgb = SF_GRAY
p_date.font.italic = True

# ============= SLIDE 2: KEY HIGHLIGHTS =============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide2)
add_footer(slide2, 2)

# Title
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
tf2 = title2.text_frame
tf2.text = "Key Highlights"
p2 = tf2.paragraphs[0]
p2.font.size = Pt(44)
p2.font.bold = True
p2.font.color.rgb = SF_DARK_BLUE
p2.font.name = "Arial"

# Accent line under title
line2 = slide2.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(2), Inches(0.03))
line2.fill.solid()
line2.fill.fore_color.rgb = SF_LIGHT_BLUE
line2.line.fill.background()

# Content
content2 = slide2.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.5), Inches(3.5))
tf2c = content2.text_frame
tf2c.word_wrap = True

highlights = [
    ("Revenue increased 42% YoY to $8.7M", SF_DARK_BLUE),
    ("Customer acquisition up 67% - added 234 new customers", SF_MED_BLUE),
    ("Product launch exceeded targets by 28%", SF_TEAL),
    ("Team expansion: 12 new hires across engineering and sales", SF_ORANGE),
    ("Net Promoter Score improved from 68 to 79", SF_PINK)
]

for i, (text, color) in enumerate(highlights):
    if i == 0:
        tf2c.text = "• " + text
        p = tf2c.paragraphs[0]
    else:
        p = tf2c.add_paragraph()
        p.text = "• " + text
    p.font.size = Pt(20)
    p.font.color.rgb = color
    p.space_before = Pt(14)
    p.font.name = "Arial"

# ============= SLIDE 3: REVENUE BREAKDOWN =============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide3)
add_footer(slide3, 3)

title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
tf3 = title3.text_frame
tf3.text = "Revenue Breakdown"
p3 = tf3.paragraphs[0]
p3.font.size = Pt(44)
p3.font.bold = True
p3.font.color.rgb = SF_DARK_BLUE
p3.font.name = "Arial"

line3 = slide3.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(2.8), Inches(0.03))
line3.fill.solid()
line3.fill.fore_color.rgb = SF_LIGHT_BLUE
line3.line.fill.background()

# Total Revenue callout box
callout = slide3.shapes.add_shape(1, Inches(6.5), Inches(1.2), Inches(3), Inches(1))
callout.fill.solid()
callout.fill.fore_color.rgb = SF_MED_BLUE
callout.line.fill.background()

callout_text = slide3.shapes.add_textbox(Inches(6.7), Inches(1.4), Inches(2.6), Inches(0.6))
tf_call = callout_text.text_frame
tf_call.text = "$8.7M\n+42% YoY"
for para in tf_call.paragraphs:
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.CENTER

# Left column - Product Line
content3a = slide3.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(4.5), Inches(3))
tf3a = content3a.text_frame

tf3a.text = "By Product Line"
p = tf3a.paragraphs[0]
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = SF_DARK_BLUE
p.space_after = Pt(12)

lines = [
    "Enterprise: $5.2M (60%)",
    "SMB: $2.4M (28%)",
    "Self-Service: $1.1M (12%)"
]
for line in lines:
    p = tf3a.add_paragraph()
    p.text = "  • " + line
    p.font.size = Pt(18)
    p.font.color.rgb = SF_GRAY
    p.space_before = Pt(8)
    p.font.name = "Arial"

# Right column - Region
content3b = slide3.shapes.add_textbox(Inches(0.7), Inches(3), Inches(4.5), Inches(2))
tf3b = content3b.text_frame

tf3b.text = "By Region"
p = tf3b.paragraphs[0]
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = SF_DARK_BLUE
p.space_after = Pt(12)

regions = [
    "North America: $6.1M (70%)",
    "EMEA: $1.7M (20%)",
    "APAC: $0.9M (10%)"
]
for reg in regions:
    p = tf3b.add_paragraph()
    p.text = "  • " + reg
    p.font.size = Pt(18)
    p.font.color.rgb = SF_GRAY
    p.space_before = Pt(8)
    p.font.name = "Arial"

# ============= SLIDE 4: CUSTOMER GROWTH =============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide4)
add_footer(slide4, 4)

title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
tf4 = title4.text_frame
tf4.text = "Customer Growth & Retention"
p4 = tf4.paragraphs[0]
p4.font.size = Pt(44)
p4.font.bold = True
p4.font.color.rgb = SF_DARK_BLUE
p4.font.name = "Arial"

line4 = slide4.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(4), Inches(0.03))
line4.fill.solid()
line4.fill.fore_color.rgb = SF_LIGHT_BLUE
line4.line.fill.background()

# Key metrics boxes
metrics = [
    ("234", "New Customers", SF_MED_BLUE),
    ("1,456", "Total Active", SF_TEAL),
    ("2.3%", "Churn Rate", SF_PINK)
]

x_pos = 0.5
for number, label, color in metrics:
    box = slide4.shapes.add_shape(1, Inches(x_pos), Inches(1.2), Inches(2.8), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    num_box = slide4.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.35), Inches(2.4), Inches(0.4))
    tfn = num_box.text_frame
    tfn.text = number
    tfn.paragraphs[0].font.size = Pt(36)
    tfn.paragraphs[0].font.bold = True
    tfn.paragraphs[0].font.color.rgb = WHITE
    tfn.paragraphs[0].alignment = PP_ALIGN.CENTER

    lab_box = slide4.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.75), Inches(2.4), Inches(0.25))
    tfl = lab_box.text_frame
    tfl.text = label
    tfl.paragraphs[0].font.size = Pt(14)
    tfl.paragraphs[0].font.color.rgb = WHITE
    tfl.paragraphs[0].alignment = PP_ALIGN.CENTER

    x_pos += 3.1

# Additional details
content4 = slide4.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(8.5), Inches(2.5))
tf4c = content4.text_frame

details = [
    "Customer Lifetime Value: $47,300",
    "Average Deal Size: $37,200 (+23% QoQ)",
    "",
    "Top 3 Customer Wins:",
    "  → Global Financial Services - $420K ARR",
    "  → Healthcare Platform - $280K ARR",
    "  → E-commerce Leader - $195K ARR"
]

for i, line in enumerate(details):
    if i == 0:
        tf4c.text = line
        p = tf4c.paragraphs[0]
    else:
        p = tf4c.add_paragraph()
        p.text = line

    if line.startswith("Top"):
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = SF_DARK_BLUE
        p.space_before = Pt(16)
    elif line == "":
        continue
    else:
        p.font.size = Pt(16)
        p.font.color.rgb = SF_GRAY
        p.space_before = Pt(8)
    p.font.name = "Arial"

# ============= SLIDE 5: Q2 OUTLOOK =============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide5)
add_footer(slide5, 5)

title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
tf5 = title5.text_frame
tf5.text = "Q2 2026 Outlook"
p5 = tf5.paragraphs[0]
p5.font.size = Pt(44)
p5.font.bold = True
p5.font.color.rgb = SF_DARK_BLUE
p5.font.name = "Arial"

line5 = slide5.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(2.3), Inches(0.03))
line5.fill.solid()
line5.fill.fore_color.rgb = SF_LIGHT_BLUE
line5.line.fill.background()

# Target callout
target_box = slide5.shapes.add_shape(1, Inches(6.8), Inches(1.2), Inches(2.7), Inches(0.7))
target_box.fill.solid()
target_box.fill.fore_color.rgb = SF_ORANGE
target_box.line.fill.background()

target_text = slide5.shapes.add_textbox(Inches(7), Inches(1.35), Inches(2.3), Inches(0.4))
tft = target_text.text_frame
tft.text = "Target: $10.2M"
tft.paragraphs[0].font.size = Pt(24)
tft.paragraphs[0].font.bold = True
tft.paragraphs[0].font.color.rgb = WHITE
tft.paragraphs[0].alignment = PP_ALIGN.CENTER

# Content sections
content5 = slide5.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.5), Inches(3.6))
tf5c = content5.text_frame

sections = [
    ("Key Initiatives", [
        "Launch AI-powered analytics module (May)",
        "Expand into 3 new verticals",
        "Open APAC office in Singapore",
        "Release mobile app (beta in June)"
    ], SF_MED_BLUE),
    ("Strategic Focus", [
        "Accelerate enterprise pipeline",
        "Enhance product-market fit in healthcare",
        "Scale customer success operations"
    ], SF_TEAL)
]

first = True
for section_title, items, color in sections:
    if not first:
        p = tf5c.add_paragraph()
        p.text = ""
        p.space_before = Pt(20)

    p = tf5c.add_paragraph() if not first else tf5c.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = color
    p.space_after = Pt(10)
    p.font.name = "Arial"

    for item in items:
        p = tf5c.add_paragraph()
        p.text = "  • " + item
        p.font.size = Pt(16)
        p.font.color.rgb = SF_GRAY
        p.space_before = Pt(6)
        p.font.name = "Arial"

    first = False

# Save presentation
output_file = 'Q1_2026_Results_Salesforce_Branded.pptx'
prs.save(output_file)
print(f"✓ Salesforce-branded presentation created: {output_file}")
print(f"  - Official SF 2025 Corporate color palette")
print(f"  - Professional layout with headers and footers")
print(f"  - 5 slides with engaging data visualization")
